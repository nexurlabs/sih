#!/usr/bin/env python3
"""Fill the user-provided SIH 2026 idea template without text overflow."""
from __future__ import annotations

import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / 'docs' / 'SIH2026-IDEA-Presentation-Format.pptx'
CACHE_TEMPLATE = Path('/root/.hermes/cache/documents/doc_aee6b1a033e8_SIH2026-IDEA-Presentation-Format.pptx')
OUT = ROOT / 'docs' / 'MailTrace_SIH26106.pptx'
PDF_OUT = ROOT / 'docs' / 'MailTrace_SIH26106_idea.pdf'
SHOTS = ROOT / 'docs' / 'shots'

# The SIH template is 13.333 x 7.5 inches. Keep all new content inside the
# large body region (x=.55..12.75, y=1.35..6.65).
NAVY = RGBColor(24, 33, 43)
BLUE = RGBColor(18, 110, 170)
RUST = RGBColor(169, 69, 49)
MUTED = RGBColor(94, 104, 112)
RULE = RGBColor(203, 209, 212)
PAPER = RGBColor(251, 250, 246)
SAND = RGBColor(241, 238, 231)
PALE_BLUE = RGBColor(234, 243, 248)
PALE_RUST = RGBColor(247, 233, 228)
WHITE = RGBColor(255, 255, 255)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def remove_text_containing(slide, snippets: tuple[str, ...]) -> None:
    for shape in list(slide.shapes):
        value = getattr(shape, 'text', '') or ''
        if any(s.lower() in value.lower() for s in snippets):
            remove_shape(shape)


def set_shape_text(shape, value: str, size: int = 10, color=NAVY, bold=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = 'Arial'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, value, size: float = 12, color=NAVY, bold=False,
             font='Arial', align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, italic=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return shape


def add_rich_text(slide, x, y, w, h, lines, size=10, color=NAVY, leading=1.05,
                  bullet=False, bullet_color=RUST):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for idx, item in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        p.line_spacing = leading
        if bullet:
            p.text = item
            p.level = 0
            p.font.name = 'Arial'
            p.font.size = Pt(size)
            p.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = 'Arial'
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return shape


def add_panel(slide, x, y, w, h, fill=PAPER, line=RULE, width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(width)
    return shape


def add_label(slide, x, y, value, color=BLUE, w=2.2):
    add_text(slide, x, y, w, 0.22, value.upper(), 8, color, True, 'Arial', margin=0)


def add_connector(slide, x1, y1, x2, y2, color=MUTED, width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_flow_box(slide, x, y, w, h, number, title, detail, accent=BLUE, fill=PALE_BLUE):
    add_panel(slide, x, y, w, h, fill, RULE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    add_text(slide, x + 0.10, y + 0.14, w - 0.20, 0.18, number, 8, accent, True, 'Arial', margin=0)
    add_text(slide, x + 0.10, y + 0.38, w - 0.20, 0.24, title, 10, NAVY, True, 'Arial', margin=0)
    add_text(slide, x + 0.10, y + 0.68, w - 0.20, h - 0.78, detail, 8.1, MUTED, False, 'Arial', margin=0)


def add_bullets(slide, x, y, w, h, items, size=10, color=NAVY):
    # Draw small rust dots ourselves for a stable, compact list.
    yy = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.08), Inches(0.07), Inches(0.07))
        dot.fill.solid(); dot.fill.fore_color.rgb = RUST; dot.line.fill.background()
        add_text(slide, x + 0.15, yy, w - 0.15, 0.36, item, size, color, False, 'Arial', margin=0)
        yy += 0.43


def add_compact_bullets(slide, x, y, w, items, size=7.6, step=0.23):
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.055), Inches(0.055))
        dot.fill.solid(); dot.fill.fore_color.rgb = RUST; dot.line.fill.background()
        add_text(slide, x + 0.13, y, w - 0.13, 0.18, item, size, NAVY, False, 'Arial', margin=0)
        y += step


def clean_template(pres: Presentation) -> None:
    # Delete the instruction page. The portal wants max six slides including
    # the title page.
    if len(pres.slides) > 6:
        slide_id = pres.slides._sldIdLst[6]
        pres.part.drop_rel(slide_id.rId)
        pres.slides._sldIdLst.remove(slide_id)
    for idx, slide in enumerate(pres.slides, 1):
        remove_text_containing(slide, (
            'Problem Statement ID', 'Proposed Solution (Describe',
            'Technologies to be used', 'Analysis of the feasibility',
            'Potential impact on the target audience',
            'Details / Links of the reference',
        ))
        # The supplied master has a decorative oval that cuts through the
        # custom heading. Remove only that ornament; keep the official logos.
        for shape in list(slide.shapes):
            if getattr(shape, 'name', '').startswith('Oval'):
                remove_shape(shape)
        # Replace the small top-left team label without changing its box.
        for shape in slide.shapes:
            if getattr(shape, 'text', '').strip() == 'Your Team Name':
                set_shape_text(shape, 'TEAM NAME', 8, BLUE, True)
        # Slide numbers remain in the supplied footer.
        for shape in slide.shapes:
            if getattr(shape, 'text', '').strip() == str(idx):
                set_shape_text(shape, str(idx), 9, WHITE, False)


def title_slide(slide, portal_title: str, team_id: str, team_name: str):
    remove_text_containing(slide, ('TITLE PAGE', 'SMART INDIA HACKATHON 2026'))
    add_text(slide, 0.55, 0.55, 6.3, 0.38, 'SMART INDIA HACKATHON 2026', 18, BLUE, True, 'Georgia', margin=0)
    add_text(slide, 0.55, 1.05, 5.95, 0.52, 'MailTrace', 30, NAVY, True, 'Georgia', margin=0)
    add_text(slide, 0.55, 1.58, 5.95, 0.36, 'Explainable email threat detection with origin and campaign forensics', 12.5, NAVY, True, 'Arial', margin=0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.98), Inches(0.9), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = RUST; bar.line.fill.background()
    add_text(slide, 0.55, 2.13, 5.95, 0.42, 'Upload a saved email. See the evidence behind its risk and likely infrastructure.', 10.5, MUTED, False, 'Arial', margin=0)
    # Keep all required title-page fields in one bounded left-hand column.
    fields = [
        ('PROBLEM STATEMENT ID', 'SIH26106'),
        ('PROBLEM STATEMENT TITLE', portal_title),
        ('THEME', 'Blockchain & Cybersecurity'),
        ('PS CATEGORY', 'Software'),
        ('TEAM ID', team_id),
        ('TEAM NAME', team_name),
    ]
    y = 2.72
    for key, value in fields:
        add_text(slide, 0.55, y, 1.58, 0.16, key, 7.2, MUTED, True, 'Arial', margin=0)
        add_text(slide, 2.18, y - 0.02, 4.22, 0.24 if key != 'PROBLEM STATEMENT TITLE' else 0.42, value, 8.5 if key != 'PROBLEM STATEMENT TITLE' else 7.9, RUST if 'ENTER' in value else NAVY, 'ENTER' not in value, 'Arial', margin=0)
        y += 0.47 if key != 'PROBLEM STATEMENT TITLE' else 0.64


def slide2(slide):
    add_text(slide, 0.55, 1.08, 12.0, 0.20, 'From one suspicious message to a connected case file.', 9.5, NAVY, True, 'Arial', margin=0)
    add_label(slide, 0.55, 1.38, 'Proposed solution')
    xs = [0.55, 3.45, 6.35, 9.25]
    content = [
        ('01', 'CAPTURE', '.eml bytes stay local; hash exact evidence.'),
        ('02', 'EXTRACT', 'Headers, links, MIME parts and Received chain.'),
        ('03', 'EXPLAIN', 'Auth evidence, score and citeable reasons.'),
        ('04', 'CONNECT', 'Shared Reply-To, domain or hop IP becomes a campaign lead.'),
    ]
    for i, (x, (n, t, d)) in enumerate(zip(xs, content)):
        add_flow_box(slide, x, 1.70, 2.48, 1.20, n, t, d, RUST if i == 2 else BLUE, PALE_RUST if i == 2 else PALE_BLUE)
        if i < 3:
            add_connector(slide, x + 2.48, 2.30, x + 2.82, 2.30, MUTED, 1.4)
    add_label(slide, 0.55, 3.10, 'Innovation proof', RUST, 2.4)
    add_panel(slide, 0.55, 3.38, 12.15, 2.48, PAPER, RULE)
    add_text(slide, 0.78, 3.57, 12.0, 0.22, 'CAMPAIGN CORRELATION  /  TWO CASES, ONE EXPLAINABLE RELATIONSHIP', 8, BLUE, True, 'Arial', margin=0)
    # Large two-node diagram. It is intentionally sparse for projector scale.
    add_panel(slide, 0.95, 4.13, 3.65, 1.18, SAND, NAVY)
    add_text(slide, 1.16, 4.30, 3.2, 0.20, '07  /  CLOUD HOPS', 8.5, BLUE, True, 'Arial', margin=0)
    add_text(slide, 1.16, 4.62, 3.2, 0.20, 'mail-secure.net', 13, NAVY, True, 'Arial', margin=0)
    add_text(slide, 1.16, 4.96, 3.2, 0.18, 'Reply-To + 18.184.10.20 / AWS', 7.5, MUTED, margin=0)
    add_panel(slide, 8.65, 4.13, 3.65, 1.18, PALE_RUST, RUST)
    add_text(slide, 8.86, 4.30, 3.2, 0.20, '08  /  CAMPAIGN TWIN', 8.5, RUST, True, 'Arial', margin=0)
    add_text(slide, 8.86, 4.62, 3.2, 0.20, 'mail-secure.net', 13, NAVY, True, 'Arial', margin=0)
    add_text(slide, 8.86, 4.96, 3.2, 0.18, 'Reply-To + 18.184.10.20 / AWS', 7.5, MUTED, margin=0)
    add_connector(slide, 4.60, 4.72, 8.65, 4.72, RUST, 2.5)
    add_panel(slide, 5.63, 4.38, 2.00, 0.70, PAPER, RUST)
    add_text(slide, 5.79, 4.51, 1.68, 0.15, 'SHARED', 8, RUST, True, 'Arial', PP_ALIGN.CENTER, margin=0)
    add_text(slide, 5.79, 4.72, 1.68, 0.20, 'Reply-To + hop + domain', 7.4, NAVY, False, 'Arial', PP_ALIGN.CENTER, margin=0)
    add_text(slide, 0.78, 5.62, 11.6, 0.16, 'Correlation creates a campaign candidate — never an identity or exact-location claim.', 8.3, MUTED, False, 'Arial', margin=0)


def slide3(slide):
    add_text(slide, 0.55, 1.08, 12.0, 0.20, 'Evidence pipeline: small, explainable modules.', 9.5, NAVY, True, 'Arial', margin=0)
    add_label(slide, 0.55, 1.38, 'Methodology')
    modules = [
        ('01  .EML', 'raw bytes + SHA-256', BLUE),
        ('02  MIME', 'headers, body, URLs', BLUE),
        ('03  AUTH', 'SPF / DKIM / DMARC', RUST),
        ('04  ORIGIN', 'hops + hosting intel', BLUE),
        ('05  OUTPUT', 'score + map + graph + PDF', RUST),
    ]
    x = 0.55
    for i, (title, detail, accent) in enumerate(modules):
        add_flow_box(slide, x, 1.70, 2.35, 1.10, title[:2], title[3:], detail, accent, PALE_RUST if accent == RUST else PALE_BLUE)
        if i < 4:
            add_connector(slide, x + 2.35, 2.25, x + 2.63, 2.25, MUTED, 1.3)
        x += 2.55
    add_label(slide, 0.55, 3.08, 'Implemented prototype', GREEN if 'GREEN' in globals() else BLUE, 2.8)
    add_bullets(slide, 0.55, 3.38, 5.45, 2.0, [
        'FastAPI + Python email parser + SQLite case store',
        'NetworkX campaign graph + Leaflet hop map',
        'ReportLab forensic PDF with exact-file SHA-256',
        'Offline demo intelligence keeps laptop runs reproducible',
    ], 9.2)
    add_label(slide, 6.35, 3.08, 'Real localhost proof', RUST, 2.8)
    if (SHOTS / 'case_header.png').exists():
        slide.shapes.add_picture(str(SHOTS / 'case_header.png'), Inches(6.35), Inches(3.38), width=Inches(6.35), height=Inches(0.92))
    if (SHOTS / 'case_lower.png').exists():
        slide.shapes.add_picture(str(SHOTS / 'case_lower.png'), Inches(6.35), Inches(4.40), width=Inches(6.35), height=Inches(1.13))
    add_text(slide, 6.35, 5.65, 6.2, 0.18, 'Captured from the local case screen: score, reasons, hop map and evidence footer.', 7.5, MUTED, margin=0)


def slide4(slide):
    add_text(slide, 0.55, 1.08, 12.0, 0.20, 'Feasible now; harden the edges as adoption grows.', 9.5, NAVY, True, 'Arial', margin=0)
    add_label(slide, 0.55, 1.38, 'Feasibility')
    add_panel(slide, 0.55, 1.72, 5.85, 2.28, PALE_BLUE, RULE)
    add_text(slide, 0.80, 1.96, 5.3, 0.23, 'NOW  /  IDEA-ROUND PROTOTYPE', 10, BLUE, True, 'Arial', margin=0)
    add_bullets(slide, 0.82, 2.30, 5.2, 1.45, [
        'Saved .eml upload; no mailbox credentials',
        'Eight crafted cases for spoof, phish, BEC and campaign paths',
        'Offline demo geo/domain cache for reproducible runs',
        'Local case store, graph, map and forensic PDF',
    ], 9.1)
    add_panel(slide, 6.78, 1.72, 5.92, 2.28, PALE_RUST, RULE)
    add_text(slide, 7.03, 1.96, 5.35, 0.23, 'NEXT  /  PRODUCTION HARDENING', 10, RUST, True, 'Arial', margin=0)
    add_bullets(slide, 7.05, 2.30, 5.25, 1.45, [
        'Live DNS and cryptographic auth verification',
        'Maintained GeoIP/ASN data and private deployment',
        'Access control, PII masking and retention policies',
        'Larger labelled corpus, calibration and analyst feedback',
    ], 9.1)
    add_label(slide, 0.55, 4.27, 'Risks and mitigations', RUST, 2.6)
    add_panel(slide, 0.55, 4.58, 12.15, 1.28, PAPER, RULE)
    add_text(slide, 0.78, 4.75, 2.3, 0.18, 'RISK', 8, BLUE, True, 'Arial', margin=0)
    add_text(slide, 3.05, 4.75, 9.2, 0.18, 'MITIGATION', 8, BLUE, True, 'Arial', margin=0)
    rows = [
        ('Headers may be forged', 'Show uncertainty and raw evidence; never call a lead proof.'),
        ('IP geo is not a person', 'Report hosting city/ISP only; explicitly reject GPS attribution.'),
        ('Email is sensitive', 'Local-first workflow; no Gmail ingest; future masking and retention controls.'),
    ]
    yy = 5.02
    for i, (a, b) in enumerate(rows):
        if i:
            add_connector(slide, 0.78, yy - 0.08, 12.35, yy - 0.08, RULE, 0.6)
        add_text(slide, 0.78, yy, 2.1, 0.25, a, 8.4, MUTED, True, 'Arial', margin=0)
        add_text(slide, 3.05, yy, 9.1, 0.25, b, 8.4, NAVY, margin=0)
        yy += 0.28


def slide5(slide):
    add_text(slide, 0.55, 1.08, 12.0, 0.20, 'Impact: turn a vague warning into an actionable case.', 9.5, NAVY, True, 'Arial', margin=0)
    add_label(slide, 0.55, 1.38, 'Analyst difference')
    add_panel(slide, 0.55, 1.72, 3.35, 1.48, PALE_BLUE, RULE)
    add_text(slide, 0.78, 1.98, 2.8, 0.18, 'FILTER', 9, BLUE, True, 'Arial', margin=0)
    add_text(slide, 0.78, 2.32, 2.8, 0.30, 'spam / not spam', 19, NAVY, True, 'Georgia', margin=0)
    add_text(slide, 0.78, 2.76, 2.8, 0.20, 'A label leaves the next question unanswered.', 8.2, MUTED, margin=0)
    add_connector(slide, 4.05, 2.47, 4.62, 2.47, RUST, 2.5)
    add_panel(slide, 4.88, 1.62, 7.82, 1.68, PALE_RUST, RUST)
    add_text(slide, 5.15, 1.90, 7.15, 0.18, 'MAILTRACE CASE FILE', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 5.15, 2.25, 7.15, 0.25, 'SPF/DKIM/DMARC fail  ·  header evidence  ·  earliest hop', 11.5, NAVY, True, 'Arial', margin=0)
    add_text(slide, 5.15, 2.62, 7.15, 0.20, 'Frankfurt / AWS  ·  related case 08  ·  SHA-256', 9.2, MUTED, margin=0)
    add_label(slide, 0.55, 3.63, 'Target users and value', BLUE, 2.8)
    add_panel(slide, 0.55, 3.92, 12.15, 1.22, PAPER, RULE)
    rows = [
        ('IT / helpdesk', 'Triage a suspicious message with reasons they can cite.'),
        ('SOC / incident response', 'Preserve original evidence and connect related messages.'),
        ('Universities / SMEs', 'Run locally without handing an inbox to a third party.'),
    ]
    yy = 4.14
    for i, (a, b) in enumerate(rows):
        if i:
            add_connector(slide, 0.78, yy - 0.08, 12.35, yy - 0.08, RULE, 0.6)
        add_text(slide, 0.78, yy, 2.4, 0.20, a, 8.8, MUTED, True, 'Arial', margin=0)
        add_text(slide, 3.20, yy, 9.0, 0.20, b, 8.8, NAVY, margin=0)
        yy += 0.30
    add_label(slide, 0.55, 5.45, 'Pilot measures  /  no invented numbers', RUST, 3.7)
    add_text(slide, 4.25, 5.45, 8.4, 0.18, 'time-to-triage  ·  evidence completeness  ·  linked-campaign discovery  ·  analyst agreement', 7.6, MUTED, margin=0)


def slide6(slide):
    remove_text_containing(slide, ('Details / Links',))
    add_text(slide, 0.55, 1.08, 12.0, 0.20, 'Research basis and implementation references.', 9.5, NAVY, True, 'Arial', margin=0)
    add_label(slide, 0.55, 1.38, 'Problem + standards', BLUE, 2.6)
    add_panel(slide, 0.55, 1.68, 5.85, 2.65, PAPER, RULE)
    add_text(slide, 0.78, 1.92, 5.3, 0.20, 'SIH26106', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 1.98, 1.92, 4.05, 0.68, 'Target scope: NLP/ML detection, header/origin forensics, domain intelligence, graph correlation and reporting. Hybrid score: forensic header rules plus a bounded local NLP component; optional Groq Qwen3.8-27B output is redacted, advisory-only, disabled by default, and not a validated detector.', 7.4, NAVY, margin=0)
    add_text(slide, 0.78, 2.65, 5.3, 0.20, 'RFC 5322', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 1.98, 2.65, 4.05, 0.40, 'Internet Message Format: structured headers, message IDs and Received fields.', 8.5, NAVY, margin=0)
    add_text(slide, 0.78, 3.32, 5.3, 0.20, 'RFC 7208 / 6376 / 7489', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 2.45, 3.32, 3.58, 0.40, 'SPF, DKIM and DMARC concepts used for authentication evidence and alignment.', 8.5, NAVY, margin=0)
    add_label(slide, 6.78, 1.38, 'Build references', RUST, 2.2)
    add_panel(slide, 6.78, 1.68, 5.92, 2.65, PAPER, RULE)
    add_text(slide, 7.02, 1.92, 1.45, 0.20, 'Python email', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 8.50, 1.92, 3.90, 0.42, 'MIME parsing; attachments named and hashed, never executed.', 8.5, NAVY, margin=0)
    add_text(slide, 7.02, 2.55, 1.45, 0.20, 'NetworkX + Leaflet', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 8.50, 2.55, 3.90, 0.42, 'Relationship graph and hop map for analyst inspection.', 8.5, NAVY, margin=0)
    add_text(slide, 7.02, 3.18, 1.45, 0.20, 'ReportLab + SHA-256', 9, RUST, True, 'Arial', margin=0)
    add_text(slide, 8.50, 3.18, 3.90, 0.52, 'Portable case PDF and integrity fingerprint of the exact uploaded bytes.', 8.5, NAVY, margin=0)
    add_label(slide, 0.55, 4.48, 'Scope and safety', BLUE, 2.1)
    add_panel(slide, 0.55, 4.76, 12.15, 0.82, PALE_BLUE, RULE)
    add_compact_bullets(slide, 0.80, 4.91, 11.6, [
        'Local crafted .eml files only; no Gmail/Outlook ingest and no real phishing messages.',
        'Geo output is hosting infrastructure with stated uncertainty, never the exact location of a person.',
        'Evidence hash is an integrity fingerprint, not a public blockchain or cryptocurrency ledger.',
    ], 7.4, 0.22)
    add_text(slide, 0.55, 5.78, 12.0, 0.16, 'Source links: sih.gov.in/sih2026PS  ·  rfc-editor.org  ·  docs.python.org  ·  networkx.org  ·  leafletjs.com', 7.2, MUTED, margin=0)


def submission_metadata() -> tuple[str, str, str]:
    values = {
        "team id": os.environ.get("MAILTRACE_TEAM_ID", "").strip(),
        "registered team name": os.environ.get("MAILTRACE_TEAM_NAME", "").strip(),
        "portal title": os.environ.get("MAILTRACE_PORTAL_TITLE", "").strip(),
    }
    missing = [label for label, value in values.items() if not value]
    placeholders = [
        label for label, value in values.items()
        if value and ("placeholder" in value.lower() or "[enter " in value.lower())
    ]
    if missing or placeholders:
        problems = missing + [f"{label} placeholder" for label in placeholders]
        raise SystemExit(
            "refusing to generate an upload candidate: provide real "
            + ", ".join(problems)
            + " via MAILTRACE_TEAM_ID, MAILTRACE_TEAM_NAME, and MAILTRACE_PORTAL_TITLE"
        )
    return values["portal title"], values["team id"], values["registered team name"]


def export_canonical_pdf(pptx_path: Path | None = None, pdf_path: Path | None = None) -> Path:
    source = Path(pptx_path or OUT)
    target = Path(pdf_path or PDF_OUT)
    if not source.is_file():
        raise SystemExit(f"cannot export missing editable deck: {source}")
    converter = shutil.which("soffice") or shutil.which("libreoffice")
    if not converter:
        raise SystemExit("cannot export PDF: soffice/libreoffice is not installed")
    with tempfile.TemporaryDirectory(prefix="mailtrace-pdf-") as temp_dir:
        try:
            subprocess.run(
                [converter, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, str(source)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
        except (OSError, subprocess.CalledProcessError) as exc:
            raise SystemExit("LibreOffice could not export the editable deck") from exc
        converted = Path(temp_dir) / f"{source.stem}.pdf"
        if not converted.is_file():
            raise SystemExit("LibreOffice completed without producing a PDF")
        target.parent.mkdir(parents=True, exist_ok=True)
        converted.replace(target)
    return target


def main(export_pdf: bool = False):
    portal_title, team_id, team_name = submission_metadata()
    template = TEMPLATE if TEMPLATE.exists() else CACHE_TEMPLATE
    if not template.exists():
        raise SystemExit(f'missing SIH 2026 template: {TEMPLATE}')
    pres = Presentation(str(template))
    clean_template(pres)
    title_slide(pres.slides[0], portal_title, team_id, team_name)
    slide2(pres.slides[1])
    slide3(pres.slides[2])
    slide4(pres.slides[3])
    slide5(pres.slides[4])
    slide6(pres.slides[5])
    pres.core_properties.title = 'MailTrace - SIH26106 Idea Presentation'
    pres.core_properties.subject = 'SIH 2026 idea submission'
    pres.core_properties.author = 'MailTrace team'
    pres.save(str(OUT))
    print('wrote', OUT)
    if export_pdf:
        print('wrote', export_canonical_pdf())


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        '--export-pdf',
        action='store_true',
        help=f'convert the editable deck to the canonical {PDF_OUT.name}',
    )
    main(export_pdf=parser.parse_args().export_pdf)
