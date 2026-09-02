from pathlib import Path
import sys

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from mailtrace.parse import parse_eml
from mailtrace.score import fuse


def test_clean_is_not_phish():
    p = Path(__file__).resolve().parents[1] / "samples" / "01_clean.eml"
    parsed = parse_eml(p.read_bytes(), p.name)
    out = fuse(parsed)
    assert parsed["auth"]["spf"] == "pass"
    assert out["label"] == "CLEAN"
    assert out["score"] < 45


def test_spoof_and_twin_graph_keys():
    root = Path(__file__).resolve().parents[1] / "samples"
    a = parse_eml((root / "02_display_spoof.eml").read_bytes(), "02_display_spoof.eml")
    b = fuse(a)
    assert a["gmail_wearing_title"]
    assert a["auth"]["spf"] == "fail"
    assert b["label"] in {"SPOOF", "PHISH"}
    assert b["score"] >= 55
    assert a["origin"]["city"] == "Frankfurt"
    assert a["origin"]["lat"] == 50.1109


def test_campaign_graph_and_pdf(tmp_path, monkeypatch):
    from mailtrace.graph_store import build
    from mailtrace.pdf_report import write_pdf
    from mailtrace import pdf_report

    monkeypatch.setattr(pdf_report, "OUT", tmp_path)
    root = Path(__file__).resolve().parents[1] / "samples"
    cases = []
    for name in ("07_cloud_hops.eml", "08_campaign_twin.eml"):
        p = parse_eml((root / name).read_bytes(), name)
        cases.append({"id": name[:2], "parsed": p, "fusion": fuse(p)})
    g = build(cases)
    assert len(g["nodes"]) == 2
    assert g["edges"]
    path = write_pdf(cases[0])
    assert path.exists() and path.stat().st_size > 400


def test_demo_corpus_labels():
    root = Path(__file__).resolve().parents[1] / "samples"
    expected = {
        "01_clean.eml": "CLEAN",
        "02_display_spoof.eml": "SPOOF",
        "03_lookalike.eml": "PHISH",
        "04_spf_fail.eml": "SPOOF",
        "05_bec_invoice.eml": "BEC",
        "06_cred_phish.eml": "PHISH",
        "07_cloud_hops.eml": "SPOOF",
        "08_campaign_twin.eml": "SPOOF",
    }
    assert len(list(root.glob("*.eml"))) == 8
    for name, label in expected.items():
        p = parse_eml((root / name).read_bytes(), name)
        out = fuse(p)
        assert out["label"] == label
        assert len(p["sha256"]) == 64


def test_attachment_metadata_is_hashed_without_execution():
    from email.message import EmailMessage
    import hashlib

    msg = EmailMessage()
    msg["From"] = "sender@example.org"
    msg["To"] = "analyst@example.org"
    msg["Subject"] = "Attachment test"
    msg.set_content("A harmless test message.")
    payload = b"benign attachment bytes"
    msg.add_attachment(payload, maintype="application", subtype="pdf", filename="invoice.pdf")

    parsed = parse_eml(msg.as_bytes(), "attachment-test.eml")
    assert parsed["attachments"] == [{
        "filename": "invoice.pdf",
        "content_type": "application/pdf",
        "size": len(payload),
        "sha256": hashlib.sha256(payload).hexdigest(),
        "risk": "low",
    }]


def test_store_preserves_exact_uploaded_bytes_and_custody_events(tmp_path, monkeypatch):
    from mailtrace import store

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")
    data = b"From: sender@example.org\n\nEvidence bytes must not change.\n"
    parsed = parse_eml(data, "evidence.eml")
    case = {"id": "case-001", "parsed": parsed, "fusion": fuse(parsed)}

    store.save_case(case, raw_bytes=data)
    loaded = store.load_case("case-001")

    evidence = loaded["evidence"]
    assert evidence["sha256"] == parsed["sha256"]
    assert evidence["stored"] is True
    assert (tmp_path / "evidence" / f"{parsed['sha256']}.eml").read_bytes() == data
    assert [event["action"] for event in evidence["custody_events"]] == [
        "received",
        "hashed",
        "stored",
    ]


def test_parser_exposes_header_evidence_and_provenance():
    root = Path(__file__).resolve().parents[1] / "samples"
    data = (root / "04_spf_fail.eml").read_bytes()

    parsed = parse_eml(data, "04_spf_fail.eml")

    assert "Authentication-Results:" in parsed["raw_headers"]
    assert "Return-Path:" in parsed["raw_headers"]
    assert parsed["auth"]["source"] == "message-header"
    assert parsed["auth"]["verification_mode"] == "header-stated"
    assert parsed["auth"]["verified"] is False
    assert parsed["provenance"] == {
        "input": "local-upload",
        "parser": "Python email BytesParser",
        "hash_algorithm": "SHA-256",
    }
    assert parsed["uncertainty"]


def test_conflicting_auth_headers_never_resolve_to_pass():
    data = (
        b"Authentication-Results: mx.example; spf=pass; dkim=pass; dmarc=pass\n"
        b"Received-SPF: fail (example.org: 203.0.113.10 not designated)\n"
        b"From: sender@example.org\n\nBody\n"
    )

    auth = parse_eml(data, "auth-conflict.eml")["auth"]

    assert auth["spf"] == "conflict"
    assert "spf" in auth["conflicts"]
    assert auth["dkim"] == "pass"
    assert auth["dmarc"] == "pass"


def test_api_upload_returns_re_readable_raw_evidence(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")
    data = b"From: sender@example.org\nSubject: Local test\n\nBody\n"

    with TestClient(app_module.app) as client:
        response = client.post(
            "/api/analyze",
            files={"file": ("api-evidence.eml", data, "message/rfc822")},
        )

        assert response.status_code == 200
        case = response.json()
        assert case["evidence"]["stored"] is True
        evidence_response = client.get(f"/api/case/{case['id']}/evidence")

    assert evidence_response.status_code == 200
    assert evidence_response.content == data


def test_local_retention_and_masking_controls_are_real_and_non_mutating(tmp_path, monkeypatch):
    from datetime import datetime, timedelta, timezone
    from mailtrace import store
    from mailtrace.privacy import masked_case

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")
    now = datetime(2026, 1, 1, tzinfo=timezone.utc)
    data = b"From: Jane Doe <jane.doe@example.org>\n\nContact jane.doe@example.org.\n"
    parsed = parse_eml(data, "privacy.eml")
    case = {"id": "case-privacy", "parsed": parsed, "fusion": fuse(parsed)}

    store.save_case(case, raw_bytes=data, retention_days=1, now=now)
    masked = masked_case(case)

    assert masked["parsed"]["from_addr"] == "j***@example.org"
    assert "j***@example.org" in masked["parsed"]["body"]
    assert case["parsed"]["from_addr"] == "jane.doe@example.org"
    assert case["parsed"]["body"] == "Contact jane.doe@example.org.\n"
    assert store.purge_expired(now=now + timedelta(days=2)) == ["case-privacy"]
    loaded = store.load_case("case-privacy")
    assert loaded["evidence"]["stored"] is False
    assert loaded["evidence"]["custody_events"][-1]["action"] == "deleted"
    assert not (tmp_path / "evidence" / f"{parsed['sha256']}.eml").exists()


def test_fusion_exposes_structured_signals_and_no_fake_ml_claim():
    root = Path(__file__).resolve().parents[1] / "samples"
    parsed = parse_eml((root / "04_spf_fail.eml").read_bytes(), "04_spf_fail.eml")
    out = fuse(parsed)

    spf_signal = next(signal for signal in out["signals"] if signal["code"] == "spf_fail")
    assert spf_signal["points"] == 22
    assert spf_signal["source"] == "message-header"
    assert out["method"] == "hybrid-forensic-nlp"
    assert out["model"]["status"] == "local-tfidf-logreg"
    assert out["model"]["validated"] is False
    assert out["probability"] is None
    assert any("not a probability" in item.lower() for item in out["uncertainty"])


def test_graph_rejects_domain_only_and_ip_only_links():
    from mailtrace.graph_store import build

    root = Path(__file__).resolve().parents[1] / "samples"
    cases = []
    for name in ("01_clean.eml", "04_spf_fail.eml"):
        parsed = parse_eml((root / name).read_bytes(), name)
        cases.append({"id": name[:2], "parsed": parsed, "fusion": fuse(parsed)})
    domain_only = build(cases)
    assert domain_only["edges"] == []

    def case(case_id: str, domain: str) -> dict:
        parsed = {
            "filename": f"{case_id}.eml",
            "from_domain": domain,
            "reply_domain": "",
            "return_domain": domain,
            "reply_to": "",
            "origin": {"ip": "203.0.113.10"},
        }
        return {"id": case_id, "parsed": parsed, "fusion": {"score": 1}}

    ip_only = build([case("a", "alpha.example"), case("b", "beta.example")])
    assert ip_only["edges"] == []


def test_api_case_graph_is_focused_on_current_campaign(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        client.post("/api/analyze-sample/04_spf_fail.eml")
        client.post("/api/analyze-sample/07_cloud_hops.eml")
        twin = client.post("/api/analyze-sample/08_campaign_twin.eml").json()

        graph = twin["graph"]
        readback = client.get(f"/api/case/{twin['id']}").json()

    assert len(graph["nodes"]) == 2
    assert len(graph["edges"]) == 1
    assert len(readback["graph"]["nodes"]) == 2
    assert readback["graph"]["edges"] == graph["edges"]


def test_parser_reports_alignment_and_unknown_states_explicitly():
    data = b"From: sender@example.org\nReturn-Path: <bounce@example.net>\n\nBody\n"

    alignment = parse_eml(data, "alignment.eml")["alignment"]

    assert alignment == {
        "from_reply_to": "unknown",
        "from_return_path": "mismatch",
        "overall": "mismatch",
    }


def test_pdf_contains_full_evidence_parity(tmp_path, monkeypatch):
    from email.message import EmailMessage
    from pypdf import PdfReader
    from mailtrace import pdf_report
    from mailtrace.pdf_report import write_pdf

    monkeypatch.setattr(pdf_report, "OUT", tmp_path)
    msg = EmailMessage()
    msg["Return-Path"] = "<bounce@example.org>"
    msg["From"] = "Sender <sender@example.org>"
    msg["Reply-To"] = "reply@example.net"
    msg["To"] = "analyst@example.org"
    msg["Subject"] = "Evidence parity"
    msg["Message-ID"] = "<parity-001@example.org>"
    msg["Authentication-Results"] = "mx.example; spf=pass; dkim=pass; dmarc=pass"
    msg.set_content("Review https://example.net/path before responding.")
    attachment = b"safe attachment bytes"
    msg.add_attachment(attachment, maintype="application", subtype="pdf", filename="invoice.pdf")
    data = msg.as_bytes()
    parsed = parse_eml(data, "parity.eml")
    case = {"id": "parity", "parsed": parsed, "fusion": fuse(parsed), "graph": {"edges": []}}

    path = write_pdf(case)
    text = "\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)

    for anchor in (
        "Return-Path",
        "Message-ID",
        "https://example.net/path",
        "invoice.pdf",
        "application/pdf",
        parsed["attachments"][0]["sha256"],
        "message-header",
        "header-stated",
        "mismatch",
        "Uncertainty",
        "local-upload",
        "Python email BytesParser",
    ):
        assert anchor in text


def test_api_masked_case_view_does_not_change_raw_evidence(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")
    data = b"From: Jane Doe <jane.doe@example.org>\n\nContact jane.doe@example.org.\n"

    with TestClient(app_module.app) as client:
        response = client.post(
            "/api/analyze",
            files={"file": ("masked.eml", data, "message/rfc822")},
        )
        case_id = response.json()["id"]
        masked_response = client.get(f"/api/case/{case_id}/masked")
        raw_response = client.get(f"/api/case/{case_id}/evidence")

    assert masked_response.status_code == 200
    assert masked_response.json()["parsed"]["from_addr"] == "j***@example.org"
    assert "j***@example.org" in masked_response.json()["parsed"]["body"]
    assert raw_response.status_code == 200
    assert raw_response.content == data


def test_offline_origin_and_intel_report_source_and_unknown_status():
    from mailtrace.intel import lookup_domains

    root = Path(__file__).resolve().parents[1] / "samples"
    parsed = parse_eml((root / "04_spf_fail.eml").read_bytes(), "04_spf_fail.eml")
    unknown = parse_eml(
        b"Received: from relay (relay [203.0.113.10])\n"
        b"From: sender@missing.example\n\nBody\n",
        "unknown-origin.eml",
    )

    assert parsed["origin"]["source"] == "offline-demo-table"
    assert parsed["origin"]["status"] == "known"
    assert parsed["origin"]["confidence"] == "heuristic"
    assert "last public Received" in parsed["origin"]["selection_method"]
    assert unknown["origin"]["status"] == "unknown"
    assert unknown["origin"]["city"] == "unknown"
    assert lookup_domains("missing.example")[0]["status"] == "unknown"
    assert lookup_domains("missing.example")[0]["source"] == "offline-cache"


# ---------------------------------------------------------------------------
# ui/index.html dossier parity, error handling, graph captions, map fallback.
# ---------------------------------------------------------------------------


def _index_html() -> str:
    return (Path(__file__).resolve().parents[1] / "ui" / "index.html").read_text()


def test_index_html_exposes_required_dossier_sections():
    """The dossier UI must include every evidence/privacy/provenance anchor."""
    html = _index_html()
    anchors = (
        "Return-Path",
        "Message-ID",
        "URL",
        "Attachment",
        "verificationMode",
        "verification_mode",
        "authSource",
        "auth_source",
        "alignment",
        "provenance",
        "uncertainty",
        "evidenceCustody",
        "prototypeBoundary",
        "graphCaption",
        "mapFallback",
    )
    missing = [a for a in anchors if a not in html]
    assert not missing, f"index.html missing dossier anchors: {missing}"


def test_index_html_dossier_calls_render_via_dispatcher_with_error_branch():
    """upload path must check response.ok and surface an actionable error."""
    html = _index_html()
    # must check response.ok before consuming JSON
    assert "response.ok" in html or "!r.ok" in html, "empty/invalid upload path missing response.ok guard"
    # must show an error element rather than silently throw
    assert "renderError" in html or "render-error" in html or "uploadError" in html, \
        "no visible error UI for empty/invalid uploads"
    # must not call render() with the literal text 'undefined' or null
    assert ".json())" in html, "upload path is expected to consume JSON response"


def test_index_html_graph_caption_is_derived_from_api_not_hardcoded():
    """Edge caption must come from the API; never the same literal every render."""
    html = _index_html()
    # the previous hardcoded text should be removed
    assert "Reply-To + same hop" not in html, "graph caption is still hardcoded"
    assert "+ domain" not in html or "graphCaption" in html, \
        "graph caption hardcoded literal still present"
    assert "e.caption" in html or "edge.caption" in html or "captions[" in html, \
        "edge caption not read from API payload"


def test_index_html_has_map_fallback_when_leaflet_or_tiles_fail():
    html = _index_html()
    # needs a CDN-down fallback message and a deterministic offline list
    assert "tile.openstreetmap.org" in html
    assert "mapFallback" in html or "map-fallback" in html or "Map unavailable" in html \
        or "fallback" in html.lower(), \
        "no map fallback when Leaflet/OSM fails"


def test_index_html_handles_leaflet_tile_errors_with_offline_list():
    """A loaded Leaflet script must still fall back when tiles fail."""
    html = _index_html()
    assert "tileerror" in html
    assert "offlineHopList" in html
    assert "drawOfflineHopList" in html


def test_api_graph_caption_is_meaningful_for_campaign_twin(tmp_path, monkeypatch):
    """07+08 produces an API-derived caption that mentions the actual shared keys."""
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        client.post("/api/analyze-sample/07_cloud_hops.eml")
        twin = client.post("/api/analyze-sample/08_campaign_twin.eml").json()

    assert twin["graph"]["edges"], "07+08 must produce at least one focused edge"
    edge = twin["graph"]["edges"][0]
    caption = edge["caption"].lower()
    # Caption is API-derived and references one of the actual shared indicators.
    assert any(token in caption for token in ("reply-to", "same hop", "domain")), \
        f"caption is not derived from shared keys: {edge['caption']!r}"
    # shared keys are surfaced too
    assert any(k.startswith("rt:") or k.startswith("ip:") or k.startswith("dom:")
               for k in edge["shared"])


def test_api_empty_upload_returns_400_and_no_case(tmp_path, monkeypatch):
    """Backend correctly 400s on empty upload; UI must surface this as an error."""
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        empty = client.post(
            "/api/analyze",
            files={"file": ("empty.eml", b"", "message/rfc822")},
        )
        cases = client.get("/api/cases").json()

    assert empty.status_code == 400
    assert cases == []


def test_api_invalid_upload_returns_400_and_no_case(tmp_path, monkeypatch):
    """Binary/non-message input must not become a misleading empty case."""
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        invalid = client.post(
            "/api/analyze",
            files={"file": ("not-email.bin", b"\x00\x01not an email", "application/octet-stream")},
        )
        cases = client.get("/api/cases").json()

    assert invalid.status_code == 400
    assert "email" in invalid.json()["detail"].lower()
    assert cases == []


def test_api_reset_removes_raw_evidence_for_cleared_cases(tmp_path, monkeypatch):
    """The localhost demo reset must not leave referenced raw bytes behind."""
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")
    data = b"From: reset@example.org\nSubject: Reset test\n\nBody\n"

    with TestClient(app_module.app) as client:
        response = client.post(
            "/api/analyze",
            files={"file": ("reset.eml", data, "message/rfc822")},
        )
        evidence_key = response.json()["evidence"]["storage_key"]
        assert (tmp_path / "evidence" / evidence_key).is_file()
        assert client.post("/api/reset").status_code == 200

    assert not (tmp_path / "evidence" / evidence_key).exists()


def test_api_unknown_sample_does_not_pollute_state(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        bad = client.post("/api/analyze-sample/does_not_exist.eml")
        cases = client.get("/api/cases").json()

    assert bad.status_code == 404
    assert cases == []


def test_official_deck_generator_preserves_template_headings_and_requires_metadata(tmp_path, monkeypatch):
    """The editable deck must keep the official headings and fail closed."""
    from pptx import Presentation
    from scripts import build_official_ppt as deck

    template = Path(__file__).resolve().parents[1] / "docs" / "SIH2026-IDEA-Presentation-Format.pptx"
    output = tmp_path / "MailTrace_SIH26106.pptx"
    monkeypatch.setattr(deck, "TEMPLATE", template)
    monkeypatch.setattr(deck, "CACHE_TEMPLATE", tmp_path / "missing-template.pptx")
    monkeypatch.setattr(deck, "OUT", output)
    monkeypatch.delenv("MAILTRACE_TEAM_ID", raising=False)
    monkeypatch.delenv("MAILTRACE_TEAM_NAME", raising=False)
    monkeypatch.delenv("MAILTRACE_PORTAL_TITLE", raising=False)

    try:
        deck.main()
    except SystemExit as exc:
        assert "team id" in str(exc).lower()
    else:
        raise AssertionError("deck generator must not create an upload candidate without metadata")
    assert not output.exists()

    monkeypatch.setenv("MAILTRACE_TEAM_ID", "TEAM-LOCAL-001")
    monkeypatch.setenv("MAILTRACE_TEAM_NAME", "MailTrace Local Team")
    monkeypatch.setenv("MAILTRACE_PORTAL_TITLE", "Verified SIH26106 portal title")
    deck.main()
    pres = Presentation(str(output))
    assert len(pres.slides) == 6
    slide_text = [
        " ".join(getattr(shape, "text", "").split())
        for slide in pres.slides
        for shape in slide.shapes
        if getattr(shape, "text", "").strip()
    ]
    for heading in (
        "IDEA TITLE",
        "TECHNICAL APPROACH",
        "FEASIBILITY AND VIABILITY",
        "IMPACT AND BENEFITS",
        "RESEARCH AND REFERENCES",
    ):
        assert heading in slide_text
    joined = "\n".join(slide_text)
    assert "[ENTER PORTAL TEAM ID]" not in joined
    assert "[ENTER REGISTERED TEAM NAME]" not in joined
    assert "TEAM-LOCAL-001" in joined
    assert "MailTrace Local Team" in joined
    assert "Verified SIH26106 portal title" in joined


def test_official_deck_metadata_rejects_placeholder_values(monkeypatch):
    """Passing the old placeholders must remain fail-closed."""
    from scripts import build_official_ppt as deck

    monkeypatch.setenv("MAILTRACE_TEAM_ID", "[ENTER PORTAL TEAM ID]")
    monkeypatch.setenv("MAILTRACE_TEAM_NAME", "[ENTER REGISTERED TEAM NAME]")
    monkeypatch.setenv("MAILTRACE_PORTAL_TITLE", "[ENTER PORTAL TITLE]")

    try:
        deck.submission_metadata()
    except SystemExit as exc:
        assert "placeholder" in str(exc).lower()
    else:
        raise AssertionError("placeholder metadata must not be accepted")


def test_official_deck_export_has_one_canonical_pdf_path(tmp_path, monkeypatch):
    """The export helper must target the documented idea-round PDF path."""
    from scripts import build_official_ppt as deck

    pptx_path = tmp_path / "MailTrace_SIH26106.pptx"
    pdf_path = tmp_path / "MailTrace_SIH26106_idea.pdf"
    pptx_path.write_bytes(b"pptx placeholder")
    monkeypatch.setattr(deck, "OUT", pptx_path)
    monkeypatch.setattr(deck, "PDF_OUT", pdf_path)

    class Result:
        stderr = ""

    def fake_run(command, **kwargs):
        converted = Path(command[command.index("--outdir") + 1]) / "MailTrace_SIH26106.pdf"
        converted.write_bytes(b"pdf placeholder")
        return Result()

    monkeypatch.setattr(deck.subprocess, "run", fake_run)
    assert deck.export_canonical_pdf() == pdf_path
    assert pdf_path.read_bytes() == b"pdf placeholder"
    assert not (tmp_path / "MailTrace_SIH26106.pdf").exists()


def test_04_spf_fail_dossier_facts_are_exact(tmp_path, monkeypatch):
    """Pin 04's exact, honest facts: SPF/DKIM/DMARC fail, Return-Path == From,
    no Reply-To, and probable hop 18.184.10.20 cached Frankfurt/AWS."""
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        case = client.post("/api/analyze-sample/04_spf_fail.eml").json()

    p = case["parsed"]
    assert p["auth"]["spf"] == "fail"
    assert p["auth"]["dkim"] == "fail"
    assert p["auth"]["dmarc"] == "fail"
    assert "notice@pec.edu.in" in p["return_path"] and "notice@pec.edu.in" in p["from_addr"], \
        f"Return-Path must equal From for 04: rp={p['return_path']!r} from={p['from_addr']!r}"
    assert not p["reply_to"], "04 has no Reply-To"
    hops_with_ip = [h for h in p["hops"] if h.get("ip")]
    assert any(h["ip"] == "18.184.10.20" for h in hops_with_ip), "expected 18.184.10.20 in hops"
    assert p["origin"]["city"] == "Frankfurt"
    assert p["origin"]["isp"] == "AWS"
    assert p["origin"]["kind"] == "cloud"
    # label and exact facts
    assert case["fusion"]["label"] == "SPOOF"


def test_05_bec_invoice_label_and_payment_signal():
    root = Path(__file__).resolve().parents[1] / "samples"
    parsed = parse_eml((root / "05_bec_invoice.eml").read_bytes(), "05_bec_invoice.eml")
    out = fuse(parsed)
    assert out["label"] == "BEC"
    codes = [s["code"] for s in out["signals"]]
    assert "payment_redirect" in codes
    assert any("payment" in r.lower() or "bec" in r.lower() for r in out["reasons"])


def test_07_08_isolated_graph_is_two_nodes_one_edge(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        assert client.post("/api/reset").status_code == 200
        client.post("/api/analyze-sample/07_cloud_hops.eml")
        twin = client.post("/api/analyze-sample/08_campaign_twin.eml").json()
        readback = client.get(f"/api/case/{twin['id']}").json()

    g = twin["graph"]
    assert len(g["nodes"]) == 2
    assert len(g["edges"]) == 1
    # readback must mirror the focused graph
    assert len(readback["graph"]["nodes"]) == 2
    assert len(readback["graph"]["edges"]) == 1
    # unrelated fixture pairs (01/04, 04 alone) must stay 0 edges
    assert client.post("/api/reset").status_code == 200
    client.post("/api/analyze-sample/01_clean.eml")
    clean = client.post("/api/analyze-sample/04_spf_fail.eml").json()
    assert clean["graph"]["edges"] == [], "01+04 must not link on shared institutional domain"


def test_index_html_handles_missing_or_failed_map_resources():
    """Even when window.L is absent or tile fetch fails, the UI must render a fallback."""
    html = _index_html()
    # either guards window.L or replaces the map div deterministically
    assert ("!window.L" in html or "typeof L ===" in html) or "mapFallback" in html
    # explicit offline hop list must exist so the map is still informative
    assert "offlineHopList" in html or "hop-list" in html or "hops (offline)" in html.lower() \
        or "Offline origin list" in html or "Offline hop list" in html \
        or "hops (deterministic)" in html.lower(), \
        "no offline deterministic hop list shown when map fails"


def test_official_deck_source_keeps_claims_and_04_facts_honest():
    source = (Path(__file__).resolve().parents[1] / "scripts" / "build_official_ppt.py").read_text()

    assert "Show confidence" not in source
    assert "with confidence" not in source
    assert "Reply-To mismatch" not in source
    assert "Show uncertainty and raw evidence" in source
    assert "SPF/DKIM/DMARC fail  ·  header evidence  ·  earliest hop" in source
    assert "Hybrid score: forensic header rules plus a bounded local NLP component; optional Groq Qwen3.8-27B output is redacted, advisory-only, disabled by default, and not a validated detector." in source


def test_qwen_assist_is_disabled_without_explicit_opt_in(monkeypatch):
    from mailtrace.llm_assist import build_assist

    monkeypatch.setenv("MAILTRACE_LLM_ENABLED", "0")
    monkeypatch.delenv("GROQ_API_KEY", raising=False)
    result = build_assist({"subject": "test", "body": "hello"}, {"score": 8, "label": "CLEAN"})

    assert result["status"] == "disabled"
    assert result["validated"] is False
    assert result["model"] == "qwen/qwen3.8-27b"


def test_qwen_prompt_is_redacted_and_excludes_raw_email_and_attachment_bytes():
    import json
    from mailtrace.llm_assist import _evidence_summary

    parsed = {
        "subject": "Reset jane.doe@example.org password",
        "from_addr": "jane.doe@example.org",
        "from_domain": "example.org",
        "body": "Call +91 98765 43210. bearer: super-secret-value",
        "urls": ["https://example.org/login?token=secret"],
        "attachments": [{"filename": "invoice.pdf", "content_type": "application/pdf", "size": 12, "sha256": "bytes"}],
        "auth": {"spf": "fail", "dkim": "unknown", "dmarc": "fail", "source": "message-header", "verification_mode": "header-stated"},
        "origin": {"ip": "18.184.10.20", "city": "Frankfurt", "isp": "AWS"},
        "hops": [{"ip": "18.184.10.20", "city": "Frankfurt", "isp": "AWS", "kind": "cloud", "status": "known"}],
    }

    text = json.dumps(_evidence_summary(parsed), ensure_ascii=False)
    assert "jane.doe@example.org" not in text
    assert "[EMAIL]" in text
    assert "+91 98765 43210" not in text
    assert "[PHONE]" in text
    assert "super-secret-value" not in text
    assert "token=secret" not in text
    assert "invoice.pdf" in text and "bytes" not in text
    assert "raw_headers" not in text


def test_qwen_assist_is_advisory_and_uses_exact_model_without_changing_score(monkeypatch):
    import json
    from mailtrace import llm_assist

    root = Path(__file__).resolve().parents[1] / "samples"
    parsed = parse_eml((root / "04_spf_fail.eml").read_bytes(), "04_spf_fail.eml")
    deterministic = fuse(parsed)
    captured = {}

    monkeypatch.setenv("MAILTRACE_LLM_ENABLED", "1")
    monkeypatch.setenv("GROQ_API_KEY", "local-test-key")

    def fake_post(url, payload, api_key, timeout):
        captured.update({"url": url, "payload": payload, "api_key": api_key, "timeout": timeout})
        return json.dumps({
            "choices": [{"message": {"content": json.dumps({
                "threat_types": ["spoofing"],
                "observations": ["Header-stated authentication fails."],
                "recommended_actions": ["Review raw headers manually."],
                "analyst_note": "Advisory only.",
                "needs_manual_review": True,
            })}}]
        }).encode()

    monkeypatch.setattr(llm_assist, "_http_post", fake_post)
    result = llm_assist.build_assist(parsed, deterministic)

    assert result["status"] == "available"
    assert result["model"] == "qwen/qwen3.8-27b"
    assert result["validated"] is False
    assert result["threat_types"] == ["spoofing"]
    assert captured["url"].endswith("/chat/completions")
    assert captured["payload"]["model"] == "qwen/qwen3.8-27b"
    prompt = captured["payload"]["messages"][1]["content"]
    assert "notice@pec.edu.in" not in prompt
    assert "<redacted_email_evidence>" in prompt
    assert deterministic["forensic_score"] == 64
    assert deterministic["label"] == "SPOOF"


def test_public_demo_is_read_only_and_keeps_campaign_graph(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from mailtrace import store
    import app as app_module

    monkeypatch.setattr(store, "DB", tmp_path / "cases.db")
    monkeypatch.setattr(store, "EVIDENCE_DIR", tmp_path / "evidence")

    with TestClient(app_module.app) as client:
        response = client.get("/api/demo/08_campaign_twin.eml")
        assert response.status_code == 200
        case = response.json()
        assert case["view"] == "public-demo"
        assert case["fusion"]["forensic_score"] == 66
        assert case["fusion"]["label"] == "SPOOF"
        assert case["fusion"]["llm_assist"]["status"] == "disabled"
        assert len(case["graph"]["nodes"]) == 2
        assert len(case["graph"]["edges"]) == 1
        assert client.get("/api/cases").json() == []
